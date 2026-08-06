#!/usr/bin/env python3
"""生成「打电话找客户思路话术」PPT — 简洁高级版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色 (苹果风格) ==========
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x1D, 0x1D, 0x1F)
GRAY     = RGBColor(0x86, 0x86, 0x8B)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF7)
BLUE     = RGBColor(0x00, 0x7A, 0xFF)
ORANGE   = RGBColor(0xFF, 0x95, 0x00)
GREEN    = RGBColor(0x34, 0xC7, 0x59)
RED      = RGBColor(0xFF, 0x3B, 0x30)
DARK_BLUE = RGBColor(0x00, 0x41, 0x75)

# ========== 辅助函数 ==========
def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_round_rect(slide, left, top, width, height, color, text='', font_size=14, font_color=BLACK, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
    return shape

def text(slide, left, top, width, height, txt, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT, name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tf

def multi_text(slide, left, top, width, height, lines, default_size=14, default_color=BLACK):
    """lines: [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        t, s, c, b = item[0], item[1] if len(item) > 1 else default_size, item[2] if len(item) > 2 else default_color, item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(s)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return tf

def card(slide, left, top, width, height, icon, title, items, title_color=BLUE, bg_color=WHITE):
    """一个信息卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xE5, 0xE5, 0xEA)
    shape.line.width = Pt(0.5)
    # 图标+标题
    text(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, f'{icon}  {title}', size=15, color=title_color, bold=True)
    # 内容
    body = '\n'.join([f'• {item}' for item in items])
    text(slide, left + 0.2, top + 0.6, width - 0.4, height - 0.8, body, size=12, color=GRAY)

def dialog_box(slide, left, top, width, height, label, content, label_color=RED, bg_color=RGBColor(0xF8, 0xF8, 0xFA)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xE5, 0xE5, 0xEA)
    shape.line.width = Pt(0.5)
    text(slide, left + 0.15, top + 0.1, width, 0.3, label, size=11, color=label_color, bold=True)
    text(slide, left + 0.15, top + 0.35, width - 0.3, height - 0.45, content, size=12, color=BLACK)

def step_bar(slide, left, top, steps):
    """横向步骤条 steps: [(num, title, desc), ...]"""
    w = 11.5 / len(steps)
    for i, (num, title, desc) in enumerate(steps):
        x = left + i * w
        # 竖线
        if i > 0:
            add_rect(slide, x, top + 0.15, 0.02, 0.2, RGBColor(0xD1, 0xD1, 0xD6))
        text(slide, x, top, w - 0.1, 0.25, num, size=22, color=BLUE, bold=True)
        text(slide, x, top + 0.3, w - 0.1, 0.25, title, size=12, color=BLACK, bold=True)
        text(slide, x, top + 0.55, w - 0.1, 0.4, desc, size=10, color=GRAY)

# =====================================================
# 第1页: 封面
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.06, BLUE)  # 顶部蓝线
text(slide, 1.5, 2.0, 10, 1.0, '打电话找客户', size=56, color=BLACK, bold=True)
text(slide, 1.5, 2.9, 10, 0.8, '思路与话术完全指南', size=30, color=GRAY)
add_rect(slide, 1.5, 3.7, 1.5, 0.04, BLUE)
text(slide, 1.5, 4.0, 10, 0.5, '时代风科技集团 · TikTok 官方认证服务商', size=14, color=GRAY)
text(slide, 1.5, 5.8, 10, 0.4, 'B2B 开发客户最直接的方式 — 绕过邮件大海，直接找到能做决定的人', size=13, color=GRAY)

# =====================================================
# 第2页: 为什么打电话
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 6, 0.5, '为什么 B2B 还要打电话？', size=30, color=BLACK, bold=True)

compare = [
    ('📧', '开发信', '每天发100封', '打开率不到5%', '回复率不到1%'),
    ('💬', 'WhatsApp', '随时聊', '客户已读不回只能等', '被动等待'),
    ('📞', '电话', '直接对话决策人', '3分钟出结果', '立竿见影 ✅'),
    ('🤝', '展会', '面对面', '一年几次成本高', '机会有限'),
]
for i, (icon, method, c1, c2, c3) in enumerate(compare):
    y = 1.5 + i * 1.3
    text(slide, 0.8, y, 0.6, 0.5, icon, size=28)
    text(slide, 1.5, y, 2.0, 0.5, method, size=18, color=BLACK, bold=True)
    text(slide, 4.0, y, 2.5, 0.5, c1, size=14, color=GRAY)
    text(slide, 6.5, y, 2.5, 0.5, c2, size=14, color=GRAY)
    text(slide, 9.0, y, 2.5, 0.5, c3, size=14, color=GRAY if '✅' not in c3 else BLUE)

# 右侧总结
add_round_rect(slide, 9.0, 1.5, 3.8, 4.5, LGRAY)
text(slide, 9.3, 1.7, 3.2, 0.4, '💡 一句话总结', size=14, color=BLUE, bold=True)
text(slide, 9.3, 2.2, 3.2, 1.5, '开发信发 100 封\n才有 1 个回复\n\n电话打 10 通\n可能成交 1 个', size=16, color=BLACK)
text(slide, 9.3, 4.5, 3.2, 1.0, '快 · 准 · 省\n电话是最被低估的\nB2B 开发方式', size=13, color=GRAY)

# =====================================================
# 第3页: 打电话前准备
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '打电话前，三件事必须准备好', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '没有准备的电话 = 浪费一次宝贵的首次接触机会', size=14, color=GRAY)

card(slide, 0.8, 1.8, 3.7, 4.5, '📋', '了解对方',
    ['公司叫什么、做什么产品',
     '对方什么职位？（采购还是老板）',
     '现在从哪采购？有中国供应商吗？',
     '最近有什么动态？（展会、新品）',
     '看一眼 LinkedIn 或官网，30秒就够'], BLUE)

card(slide, 4.8, 1.8, 3.7, 4.5, '📝', '准备好自己',
    ['我是谁 — 一句话讲清楚',
     '我能带来什么价值 — 不是"我们很好"，是"能帮你省多少"',
     '我想得到什么 — 样品单？邮箱？视频会议？',
     '对方常见拒绝有哪些，想好怎么回',
     '提前演练一遍，不要念稿'], ORANGE)

card(slide, 8.8, 1.8, 3.7, 4.5, '🎯', '调整心态',
    ['被拒绝很正常 — 10通电话有1通愿聊就是胜利',
     '你是帮忙的人，不是来求人的',
     '不用一次成交 — 拿到下一步就算赢',
     '像聊天一样，不要像推销',
     '微笑会影响你的声音'], GREEN)

# =====================================================
# 第4页: 开场白话术
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '电话开场白 — 30 秒定生死', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '对方接起电话的前 30 秒，决定了会继续听还是挂掉', size=14, color=GRAY)

# 方式一
add_round_rect(slide, 0.8, 1.7, 3.7, 4.5, LGRAY)
text(slide, 1.1, 1.9, 3.2, 0.4, '方式一  直接了当', size=15, color=BLUE, bold=True)
text(slide, 1.1, 2.4, 3.2, 0.3, '适合已发过邮件或知道对方', size=11, color=GRAY)
text(slide, 1.1, 2.9, 3.2, 3.0,
    '"Hi [Name], this is [You]\nfrom [Company].\n\nWe help Chinese factories\nreach global buyers\nthrough TikTok.\n\nI saw your company on LinkedIn.\nYou do [product], right?\n\nGot 2 minutes for a quick idea?"',
    size=13, color=BLACK)

# 方式二
add_round_rect(slide, 4.8, 1.7, 3.7, 4.5, LGRAY)
text(slide, 5.1, 1.9, 3.2, 0.4, '方式二  价值先行', size=15, color=ORANGE, bold=True)
text(slide, 5.1, 2.4, 3.2, 0.3, '适合完全陌生的客户', size=11, color=GRAY)
text(slide, 5.1, 2.9, 3.2, 3.0,
    '"Hi [Name], quick question —\n\nare you sourcing [product]\nfrom China right now?\n\nThe reason I ask:\nwe helped a factory like yours\ncut sourcing cost by 15%\nthrough TikTok live streaming.\n\nIs now a bad time?"',
    size=13, color=BLACK)

# 方式三
add_round_rect(slide, 8.8, 1.7, 3.7, 4.5, LGRAY)
text(slide, 9.1, 1.9, 3.2, 0.4, '方式三  借力打力', size=15, color=GREEN, bold=True)
text(slide, 9.1, 2.4, 3.2, 0.3, '适合有展会/行业共同背景', size=11, color=GRAY)
text(slide, 9.1, 2.9, 3.2, 3.0,
    '"Hi [Name], we actually\nhave a connection —\n\nWe both attended [Canton Fair].\n\nWe help manufacturers\nlike [知名客户名]\nget international buyers\nthrough TikTok.\n\nCurious if you have looked\ninto TikTok yet?"',
    size=13, color=BLACK)

# =====================================================
# 第5页: 通话流程
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '一通成功的电话长什么样？', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '6 步走，全程控制在 3-5 分钟。目标不是成交，是让对方给你下一步的机会', size=14, color=GRAY)

steps = [
    ('01', '开场', '你是谁 · 为什么打来\n能给我 2 分钟吗？', BLUE),
    ('02', '破冰', '提一件跟他们公司\n相关的事 · 夸一句', ORANGE),
    ('03', '试探', '现在从哪采购？\n满意吗？有想过\n试试别的渠道吗？', GREEN),
    ('04', '抛价值', '我们的客户通过 TikTok\n询盘涨了 3 倍\n你感兴趣吗？', RGBColor(0xAF, 0x52, 0xDE)),
    ('05', '处理异议', '先认同 · 再转折\n不要说"但是我们好"', RED),
    ('06', '收尾', '明确下一步\n发邮件/视频会/加 WhatsApp', DARK_BLUE),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 2.1
    add_round_rect(slide, x, 1.8, 1.9, 4.2, LGRAY)
    add_round_rect(slide, x + 0.55, 2.0, 0.8, 0.8, color, num, font_size=22, font_color=WHITE, bold=True)
    text(slide, x + 0.1, 3.0, 1.7, 0.4, title, size=16, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x + 0.1, 3.5, 1.7, 1.5, desc, size=12, color=GRAY, align=PP_ALIGN.CENTER)

text(slide, 0.8, 6.3, 11.5, 0.4, '💡 第一次通话不要试图当场成交。目标是让对方说 "Yes" 给你下一步的机会。', size=14, color=BLUE, bold=True)

# =====================================================
# 第6页: 拒绝应对话术
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '客户说 "不" 的时候，怎么接？', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '核心原则：不要反驳，先认同，再试探', size=14, color=GRAY)

dialog_box(slide, 0.8, 1.6, 5.8, 1.5, '🙅 客户说：', '"We already have suppliers."  （我们已经有供应商了）', RED)
dialog_box(slide, 7.4, 1.6, 5.2, 1.5, '✅ 你这样接：', '"That is great — means you know the market.\nI am not asking you to switch.\nJust thought a backup option might be useful.\nMind if I send you a quick intro?"', GREEN)

dialog_box(slide, 0.8, 3.4, 5.8, 1.5, '🙅 客户说：', '"Just send me an email."  （发邮件给我 — 其实是客气地拒绝）', RED)
dialog_box(slide, 7.4, 3.4, 5.2, 1.5, '✅ 你这样接：', '"Sure, happy to! One quick thing —\nwhat product category are you most interested in?\nI will make sure the email is relevant."\n（先问一个问题，让对方投入，然后再发邮件）', GREEN)

dialog_box(slide, 0.8, 5.2, 5.8, 1.5, '🙅 客户说：', '"I am not interested."  （我不感兴趣）', RED)
dialog_box(slide, 7.4, 5.2, 5.2, 1.5, '✅ 你这样接：', '"Totally understand. Just one question —\nare you buying from China at all?\nNo? → Alright, have a great day.\nYes? → Maybe worth a 30-second look?"', GREEN)

# =====================================================
# 第7页: 收尾方式
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '怎么收尾 — 拿到"下一步"才算赢', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '打完电话没有下一步 = 白打。根据对方兴趣程度，选一种收尾方式：', size=14, color=GRAY)

closings = [
    ('😊', '兴趣高', '对方问了价格、产品细节',
     '"Let me put together a proposal.\nCan we hop on a 15-min video call\nthis Thursday or Friday?"', '→  约视频会议', BLUE),
    ('🤔', '一般般', '对方没拒绝但也不热情',
     '"No worries, let me send over a quick\nintro with some case studies.\nWhat is the best email for you?"', '→  发邮件 + 加 WhatsApp', ORANGE),
    ('😐', '没兴趣', '对方明确说不需要',
     '"Completely understand.\nMind if I connect on LinkedIn?\nIf anything changes, you will\nknow where to find me."', '→  LinkedIn 连接', GREEN),
]
for i, (icon, level, desc, script, goal, color) in enumerate(closings):
    x = 0.8 + i * 4.1
    add_round_rect(slide, x, 1.6, 3.7, 4.8, LGRAY)
    text(slide, x + 0.2, 1.8, 3.3, 0.4, f'{icon}  {level}', size=18, color=color, bold=True)
    text(slide, x + 0.2, 2.3, 3.3, 0.4, desc, size=12, color=GRAY)
    add_rect(slide, x + 0.2, 2.9, 3.3, 0.005, RGBColor(0xE5, 0xE5, 0xEA))
    text(slide, x + 0.2, 3.1, 3.3, 1.8, script, size=13, color=BLACK)
    add_round_rect(slide, x + 0.2, 5.2, 3.3, 0.5, color, goal, font_size=12, font_color=WHITE, bold=True)

# =====================================================
# 第8页: 跟进
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text(slide, 0.8, 0.5, 10, 0.5, '电话后的跟进 — 大多数人死在这一步', size=30, color=BLACK, bold=True)
text(slide, 0.8, 1.0, 10, 0.4, '80% 的成交发生在第 5-12 次跟进。大部分人跟了 1-2 次就放弃了', size=14, color=GRAY)

follow_ups = [
    ('⏰', '打完立刻做', '（5分钟内）', [
        '把通话要点记下来',
        '马上发 WhatsApp / 邮件',
        '在笔记本里记下下次跟进时间',
    ], BLUE),
    ('📅', '第二天', '', [
        '检查对方有没有看你发的邮件',
        '没看的话换个时间再发',
        '用 LinkedIn 加对方好友',
    ], ORANGE),
    ('📆', '一周后', '', [
        '发一条简短 WhatsApp 跟进',
        '分享一条对方行业相关的资讯',
        '两次没回复 → 30天后再试',
    ], GREEN),
]
for i, (icon, title, sub, items, color) in enumerate(follow_ups):
    x = 0.8 + i * 4.1
    add_round_rect(slide, x, 1.5, 3.7, 4.0, LGRAY)
    text(slide, x + 0.2, 1.7, 3.3, 0.4, f'{icon}  {title}', size=16, color=color, bold=True)
    if sub:
        text(slide, x + 0.2, 2.1, 3.3, 0.3, sub, size=11, color=GRAY)
    body = '\n\n'.join([f'• {item}' for item in items])
    text(slide, x + 0.2, 2.5, 3.3, 2.5, body, size=13, color=BLACK)

# 重点提示
add_round_rect(slide, 0.8, 5.8, 11.7, 1.0, RGBColor(0x00, 0x7A, 0xFF))
text(slide, 1.2, 5.95, 10.8, 0.7, '💡 持续但不骚扰，是你甩开竞争对手最简单的方法。大部分人跟了 1-2 次就放弃了，你跟到第 5 次就赢了。', size=16, color=WHITE, bold=True)

# =====================================================
# 第9页: 总结
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.06, BLUE)
text(slide, 0.8, 1.0, 10, 0.6, '打电话找客户 — 记住这五句话就够了', size=36, color=BLACK, bold=True)
add_rect(slide, 0.8, 1.7, 1.5, 0.04, BLUE)

takeaways = [
    ('1', '准备比技巧重要', '花 3 分钟查一下对方公司，比你背 100 句话术都有用'),
    ('2', '30 秒决定生死', '说清你是谁、为什么打来、能不能给 2 分钟'),
    ('3', '你是来帮忙的', '不是来求人的 — 心态一变，声音都不一样'),
    ('4', '不要试图当场成交', '目标是拿到下一步（邮箱 / WhatsApp / 视频会议）'),
    ('5', '跟进才是真正的开始', '大部分人跟 1-2 次就放弃，你跟到第 5 次就赢了'),
]

for i, (num, title, desc) in enumerate(takeaways):
    y = 2.3 + i * 0.95
    add_round_rect(slide, 0.8, y, 0.7, 0.7, BLUE, num, font_size=22, font_color=WHITE, bold=True)
    text(slide, 1.8, y + 0.02, 4.0, 0.4, title, size=20, color=BLACK, bold=True)
    text(slide, 1.8, y + 0.4, 8.0, 0.35, desc, size=14, color=GRAY)

text(slide, 0.8, 7.0, 10, 0.4, '时代风科技集团 · TikTok 官方认证服务商', size=12, color=GRAY)

# ========== 保存 ==========
output_path = '/Users/tongtong/Desktop/tiktok外贸/打电话找客户思路话术.pptx'
prs.save(output_path)
print(f'✅ PPT 已生成: {output_path}')
